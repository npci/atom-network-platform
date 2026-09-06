# Copyright 2026 National Payments Corporation of India
# SPDX-License-Identifier: MIT

"""Unit tests for the shipment-override routing in change_dispatch.

The bug being fixed: manually-uploaded overrides for content-primary
doc types (manifest, prototype_screens, single-file xsd) were landing
in the wire `docx_b64` slot even when the file was textual (yaml,
html, xml). The partner's `/download/native` endpoint reads only
`content`, so the download returned the placeholder text instead of
the real bytes. Fix: route textual overrides for these doc types
through the content path only.
"""

from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import pytest

from app.services.change_dispatch import (
    OVERRIDE_PLACEHOLDER_TEXT,
    _extract_text_from_upload,
    _override_attach,
    _override_content,
    _TEXTUAL_OVERRIDE_MIMES,
    _uses_content_primary,
)

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _fake_row(tmp_path: Path, *, body: bytes, mime: str, name: str) -> SimpleNamespace:
    """Build a row-like object with the override_* fields _override_attach reads."""
    p = tmp_path / name
    p.write_bytes(body)
    return SimpleNamespace(
        override_path=str(p),
        override_mime_type=mime,
        override_filename=name,
        override_sha256=None,
        override_size_bytes=len(body),
        id="row-1",
    )


# ── Helper coverage ──────────────────────────────────────────────────────────


def test_placeholder_constant_exact_string():
    """Partner side hashes on this string — must not drift."""
    assert OVERRIDE_PLACEHOLDER_TEXT == (
        "Override attachment supplied by user; see attachments[] for the file."
    )


def test_content_primary_covers_manifest_prototype_xsd():
    assert _uses_content_primary("manifest")
    assert _uses_content_primary("prototype_screens")
    assert _uses_content_primary("xsd")
    assert not _uses_content_primary("tsd")
    assert not _uses_content_primary("cert_test_cases")
    assert not _uses_content_primary("product_deck")


def test_yaml_variants_are_textual():
    for m in ("application/x-yaml", "application/yaml", "text/yaml", "application/x-yml"):
        assert m in _TEXTUAL_OVERRIDE_MIMES, m


# ── Routing behaviour ────────────────────────────────────────────────────────


def test_yaml_manifest_override_skips_binary_slot(tmp_path):
    """A .yaml manifest override must NOT populate docx_b64.
    The caller's `_override_content` writes the yaml text into `content`."""
    row = _fake_row(tmp_path, body=b"foo: bar\n", mime="application/x-yaml", name="m.yaml")
    out: dict = {"doc_type": "manifest", "content": "seed"}
    handled = _override_attach(out, row, "docx", doc_type="manifest")
    assert handled is True
    # Binary slot MUST be empty — content path is authoritative here.
    assert "docx_b64" not in out
    assert "docx_filename" not in out


def test_html_prototype_override_skips_binary_slot(tmp_path):
    row = _fake_row(tmp_path, body=b"<html/>", mime="text/html", name="proto.html")
    out: dict = {"doc_type": "prototype_screens"}
    handled = _override_attach(out, row, "docx", doc_type="prototype_screens")
    assert handled is True
    assert "docx_b64" not in out


def test_single_file_xsd_override_skips_zip_slot(tmp_path):
    row = _fake_row(tmp_path, body=b'<xs:schema/>', mime="application/xml", name="one.xsd")
    out: dict = {"doc_type": "xsd"}
    handled = _override_attach(out, row, "xsd_zip", doc_type="xsd")
    assert handled is True
    assert "xsd_zip_b64" not in out
    assert "xsd_zip_filename" not in out


def test_multi_schema_zip_override_still_uses_zip_slot(tmp_path):
    row = _fake_row(tmp_path, body=b"PK\x03\x04...", mime="application/zip", name="schemas.zip")
    out: dict = {"doc_type": "xsd"}
    handled = _override_attach(out, row, "xsd_zip", doc_type="xsd")
    assert handled is True
    # ZIP is not textual — content-primary short-circuit doesn't fire.
    assert "xsd_zip_b64" in out
    assert out["xsd_zip_mime_type"] == "application/zip"


def test_pdf_tech_spec_override_uses_docx_slot(tmp_path):
    """PDF for TSD is legitimate — belongs in docx_b64 (binary path)."""
    row = _fake_row(tmp_path, body=b"%PDF-1.4\n...", mime="application/pdf", name="spec.pdf")
    out: dict = {"doc_type": "tsd"}
    handled = _override_attach(out, row, "docx", doc_type="tsd")
    assert handled is True
    # TSD is not a content-primary doc type — binary slot is populated.
    assert "docx_b64" in out
    assert out["docx_mime_type"] == "application/pdf"


def test_docx_generated_slot_for_product_note(tmp_path):
    """A docx override for a generated Word doc (e.g. product_note) also
    goes into docx_b64 — content-primary short-circuit skipped."""
    row = _fake_row(
        tmp_path, body=b"PK\x03\x04...docx...",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        name="note.docx",
    )
    out: dict = {"doc_type": "product_note"}
    handled = _override_attach(out, row, "docx", doc_type="product_note")
    assert handled is True
    assert "docx_b64" in out


# ── _override_content: content field carries textual overrides ───────────────


def test_yaml_override_content_returns_utf8_text(tmp_path):
    """The caller uses _override_content to compute the wire `content`
    field; for textual overrides it must return the file bytes as text
    (NOT the placeholder), so the partner's /download/native serves the
    real yaml.
    """
    row = _fake_row(tmp_path, body="services:\n  api: 42\n".encode(), mime="text/yaml", name="m.yaml")
    assert _override_content(row) == "services:\n  api: 42\n"


def test_pdf_override_content_returns_placeholder(tmp_path):
    """PDF has no text-extraction path (would need pypdf, a new dep) —
    stays on the placeholder + Download-to-view card behaviour."""
    row = _fake_row(tmp_path, body=b"%PDF-1.4\n...", mime="application/pdf", name="s.pdf")
    assert _override_content(row) == OVERRIDE_PLACEHOLDER_TEXT


# ── Binary text extraction (docx / xlsx / pptx) ──────────────────────────────
#
# Real python-docx / openpyxl / python-pptx round-trips — same libs the
# extractor uses — so a regression in the extractor path shows up as a broken
# test rather than a broken partner preview.


def _build_docx(tmp_path: Path) -> Path:
    from docx import Document
    p = tmp_path / "spec.docx"
    d = Document()
    d.add_heading("Tech Spec", level=1)
    d.add_paragraph("Introduction paragraph.")
    d.add_heading("API", level=2)
    d.add_paragraph("POST /foo returns bar.")
    d.save(str(p))
    return p


def _build_xlsx(tmp_path: Path) -> Path:
    from openpyxl import Workbook
    p = tmp_path / "cases.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "TC"
    ws.append(["ID", "Scenario", "Expected"])
    ws.append(["TC1", "happy path", "SUCCESS"])
    ws.append(["TC2", "timeout", "FAIL"])
    wb.save(str(p))
    return p


def _build_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    p = tmp_path / "deck.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    slide.shapes.title.text = "Overview"
    slide.placeholders[1].text = "Key point one\nKey point two"
    pres.save(str(p))
    return p


def test_docx_extractor_returns_markdown_with_headings(tmp_path):
    p = _build_docx(tmp_path)
    out = _extract_text_from_upload(str(p), _DOCX_MIME)
    assert out is not None
    assert "# Tech Spec" in out
    assert "## API" in out
    assert "Introduction paragraph." in out
    assert "POST /foo returns bar." in out


def test_docx_override_content_uses_extracted_text_not_placeholder(tmp_path):
    p = _build_docx(tmp_path)
    row = SimpleNamespace(override_path=str(p), override_mime_type=_DOCX_MIME)
    out = _override_content(row)
    # Real extracted text — NOT the sentinel — flows into wire `content`.
    assert out and out != OVERRIDE_PLACEHOLDER_TEXT
    assert "# Tech Spec" in out


def test_xlsx_extractor_returns_markdown_table_per_sheet(tmp_path):
    p = _build_xlsx(tmp_path)
    out = _extract_text_from_upload(str(p), _XLSX_MIME)
    assert out is not None
    assert "## Sheet: TC" in out
    assert "| ID | Scenario | Expected |" in out
    assert "| TC1 | happy path | SUCCESS |" in out


def test_pptx_extractor_returns_markdown_slides(tmp_path):
    p = _build_pptx(tmp_path)
    out = _extract_text_from_upload(str(p), _PPTX_MIME)
    assert out is not None
    assert "## Slide 1: Overview" in out
    assert "- Key point one" in out
    assert "- Key point two" in out


def test_docx_override_still_populates_binary_slot(tmp_path):
    """Even though `content` now carries the extracted markdown, the
    original .docx bytes MUST still ride in `docx_b64` so the partner
    Download button serves the pristine file (not the markdown text)."""
    p = _build_docx(tmp_path)
    row = SimpleNamespace(
        override_path=str(p), override_mime_type=_DOCX_MIME,
        override_filename="spec.docx", override_sha256=None,
        override_size_bytes=p.stat().st_size, id="row-x",
    )
    out: dict = {"doc_type": "tsd"}
    handled = _override_attach(out, row, "docx", doc_type="tsd")
    assert handled is True
    assert "docx_b64" in out                     # pristine binary preserved
    assert out["docx_filename"] == "spec.docx"


def _build_pdf_with_text(tmp_path: Path, text: str) -> Path:
    """Build a real PDF with visible text using reportlab if available, else
    fall back to a blank pypdf-only PDF (extractor returns None for blank)."""
    p = tmp_path / "spec.pdf"
    try:
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(p))
        c.drawString(100, 720, text)
        c.save()
    except ImportError:
        from pypdf import PdfWriter
        w = PdfWriter()
        w.add_blank_page(width=612, height=792)
        with open(p, "wb") as f:
            w.write(f)
    return p


def test_pdf_extractor_returns_text_when_pages_have_text(tmp_path):
    """PDFs with a text layer must extract to inline markdown."""
    p = _build_pdf_with_text(tmp_path, "Hello the Authority PDF world")
    out = _extract_text_from_upload(str(p), "application/pdf")
    # reportlab available → text present. reportlab absent → blank PDF → None.
    if out is None:
        pytest.skip("reportlab unavailable in this env — PDF has no text layer")
    assert "Hello the Authority PDF world" in out
    assert "## Page 1" in out


def test_pdf_blank_scanned_falls_back_to_placeholder(tmp_path):
    """A scanned / image-only PDF (no text layer) returns None → caller uses
    the Download-to-view card. OCR would need a Tesseract dep."""
    from pypdf import PdfWriter
    p = tmp_path / "scan.pdf"
    w = PdfWriter()
    w.add_blank_page(width=612, height=792)
    with open(p, "wb") as f:
        w.write(f)
    row = SimpleNamespace(override_path=str(p), override_mime_type="application/pdf")
    assert _override_content(row) == OVERRIDE_PLACEHOLDER_TEXT


def test_image_override_embeds_as_data_url_markdown(tmp_path):
    """Small images embed inline via data URL so partner ReactMarkdown
    renders them without needing a separate endpoint."""
    import base64
    p = tmp_path / "logo.png"
    # Minimal 1x1 red PNG
    p.write_bytes(base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR4nGNgAAIAAAUAAeImBZsAAAAASUVORK5CYII="
    ))
    out = _extract_text_from_upload(str(p), "image/png")
    assert out is not None
    assert out.startswith("![logo.png](data:image/png;base64,")
    assert out.endswith(")")


def test_image_override_too_large_falls_back(tmp_path):
    """Images >4 MB skip inline embed to keep payload sane — Download-to-view."""
    p = tmp_path / "big.png"
    p.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * (5 * 1024 * 1024))
    out = _extract_text_from_upload(str(p), "image/png")
    assert out is None


def test_corrupt_docx_extractor_falls_back_to_placeholder(tmp_path):
    """A file that claims docx MIME but isn't a valid docx must NOT crash
    the ship — extractor returns None → `_override_content` falls back
    to the placeholder sentinel (partner sees the Download card)."""
    p = tmp_path / "not-really.docx"
    p.write_bytes(b"this is not a real docx")
    row = SimpleNamespace(override_path=str(p), override_mime_type=_DOCX_MIME)
    assert _override_content(row) == OVERRIDE_PLACEHOLDER_TEXT
