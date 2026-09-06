# Copyright 2026 National Payments Corporation of India
# SPDX-License-Identifier: MIT

"""Product Kit packaging, publication snapshots, and partner dispatch.

Shared by the initial communication (`phase_c.communicate_change`) and the
revision flow (`negotiation_mgmt.new_version_and_ship`):

  build_kit_envelope   — assemble the exact `product_kit` A2A payload from the
                         latest version of each kit doc + XSD.
  snapshot_publication — persist an immutable KitPublication for a
                         negotiation_version (upsert; re-send of the same
                         version updates the row).
  dispatch_kit_to_partners — send the envelope over A2A. `mode="initial"`
                         dispatches ASSIGNED partners and advances them to
                         RECEIVED; `mode="revision"` re-sends to all
                         still-active partners WITHOUT touching assignment
                         status (re-acceptance is gated on the partner's
                         negotiation_version_accepted, not assignment status).
"""
import base64
import hashlib
import io
import json
import logging
import os
import zipfile
from datetime import timedelta

from sqlalchemy import select
from sqlalchemy.orm import Session

from app.core.config import settings
from app.models.base import generate_uuid, utcnow
from app.models.change_request import ChangeRequest
from app.models.kit_publication import KitPublication
from app.models.phase_c import (
    A2ATaskType, AssignmentStatus, ChangePartnerAssignment, PartnerAgent,
)
from app.models.product_kit import ProductKitDocument
from app.models.tech_spec import TechSpec
from app.models.xsd import XSD
from app.services.assignment_status import set_status
from app.services.notifications import notify_delivery_failure
from app.services.product_kit_query import latest_kit_docs

logger = logging.getLogger(__name__)


# Wire-level constant — the placeholder we put into `content` when the override
# is a binary attachment. Exported (no leading underscore alias too) so the
# partner side has a stable literal to detect "content is not the real doc".
# NEVER re-word without a matching partner change: the partner serializer
# hashes on this exact string to set `content_is_placeholder`.
OVERRIDE_PLACEHOLDER_TEXT = (
    "Override attachment supplied by user; see attachments[] for the file."
)
_OVERRIDE_PLACEHOLDER = OVERRIDE_PLACEHOLDER_TEXT  # legacy in-module alias

# Mime types we treat as human-readable text — for these the override bytes go
# into `content` (partner-side text view stays coherent). Everything else keeps
# the placeholder in `content` and ships only as an attachment. YAML variants
# added so a `.yaml` manifest override takes the content path — matches how
# the partner UI serves manifest downloads (`_NATIVE_FORMATS` in
# `app/api/dashboard/changes.py`, in the partner platform's own repository).
_TEXTUAL_OVERRIDE_MIMES = frozenset({
    "text/plain", "text/markdown", "text/csv", "text/html",
    "application/xml", "text/xml", "application/json",
    "application/x-yaml", "application/yaml", "text/yaml", "application/x-yml",
})

# doc_types whose partner-side download endpoint reads `content` (via
# `/download/native` — `_NATIVE_FORMATS` in the partner dashboard router)
# rather than `*_bytes`. A textual override for these MUST land in the wire
# `content` field and MUST NOT populate a binary slot — otherwise the
# partner UI serves the placeholder string as the "downloaded" file.
_CONTENT_PRIMARY_DOC_TYPES = frozenset({
    "manifest", "prototype_screens", "xsd",
})


def _uses_content_primary(dt: str) -> bool:
    """True when the partner-side download for this doc_type reads `content`
    (via /download/native) rather than a binary attachment. Governs whether
    an override should take the content path or the binary-slot path."""
    return (dt or "") in _CONTENT_PRIMARY_DOC_TYPES


def _extract_text_from_upload(path: str, mime: str) -> str | None:
    """Best-effort text extraction from a binary override so partner preview
    can render it inline. Returns markdown-flavoured text on success, None on
    any failure (caller falls back to the placeholder sentinel).

    Supported (all libs already in `backend/requirements.txt`, no new deps):
      .docx  → python-docx: paragraphs + Heading-N styles → `# … ## …`
      .xlsx  → openpyxl: one markdown table per sheet
      .pptx  → python-pptx: `## Slide N: <title>` + body text per slide
      .pdf   → pypdf: page text concatenated, one `## Page N` header per page
      images → base64 data-URL embedded into markdown so ReactMarkdown
               renders the image inline (partner has zero code change).

    `.mp4` / `.zip` / other unknown binaries stay on the placeholder +
    Download-to-view card behaviour. Legacy `.doc` / `.xls` / `.ppt`
    need `textract` or LibreOffice — out of scope; upload the modern
    Office XML formats or a PDF for inline preview.
    """
    m = (mime or "").lower()
    try:
        if m == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            blocks: list[str] = []
            for i, page in enumerate(reader.pages, start=1):
                try:
                    text = (page.extract_text() or "").strip()
                except Exception as exc:  # noqa: BLE001 — a corrupt page shouldn't kill the whole doc
                    logger.warning("pdf page %d extract failed path=%s: %s", i, path, exc)
                    continue
                if not text:
                    continue
                blocks.append(f"## Page {i}\n\n{text}")
            out = "\n\n".join(blocks).strip()
            # Scanned / image-only PDFs yield no text at all — signal failure
            # so the caller falls back to the Download-to-view card (OCR
            # would need a Tesseract dep, out of scope).
            return out or None

        if m.startswith("image/"):
            # Embed as a data URL so the partner's existing ReactMarkdown
            # preview renders the image inline. No partner code change,
            # no separate image endpoint. Data URLs bloat the payload a
            # bit; cap at 4 MB decoded — larger images stay as a
            # Download-to-view card.
            import base64
            size = os.path.getsize(path)
            if size > 4 * 1024 * 1024:
                logger.info("image override too large for inline embed (%d bytes) — using download card", size)
                return None
            with open(path, "rb") as f:
                data = f.read()
            b64 = base64.b64encode(data).decode()
            fname = os.path.basename(path)
            return f"![{fname}](data:{m};base64,{b64})"

        if "wordprocessingml" in m:
            from docx import Document
            doc = Document(path)
            parts: list[str] = []
            for para in doc.paragraphs:
                text = (para.text or "").strip()
                if not text:
                    parts.append("")
                    continue
                style = (getattr(para.style, "name", "") or "").lower()
                if style.startswith("heading 1") or style == "title":
                    parts.append(f"# {text}")
                elif style.startswith("heading 2"):
                    parts.append(f"## {text}")
                elif style.startswith("heading 3"):
                    parts.append(f"### {text}")
                elif style.startswith("heading "):
                    parts.append(f"#### {text}")
                else:
                    parts.append(text)
            out = "\n\n".join(p for p in parts if p is not None).strip()
            return out or None

        if "spreadsheetml" in m:
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True, read_only=True)
            blocks: list[str] = []
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                blocks.append(f"## Sheet: {ws.title}")
                header = ["" if v is None else str(v) for v in rows[0]]
                if header:
                    blocks.append("| " + " | ".join(header) + " |")
                    blocks.append("| " + " | ".join(["---"] * len(header)) + " |")
                for r in rows[1:200]:  # cap per-sheet row output
                    cells = ["" if v is None else str(v).replace("|", "\\|") for v in r]
                    if any(c.strip() for c in cells):
                        blocks.append("| " + " | ".join(cells) + " |")
                if len(rows) > 200:
                    blocks.append(f"_… {len(rows) - 200} more rows omitted from preview …_")
                blocks.append("")
            wb.close()
            out = "\n".join(blocks).strip()
            return out or None

        if "presentationml" in m:
            from pptx import Presentation
            pres = Presentation(path)
            blocks: list[str] = []
            for i, slide in enumerate(pres.slides, start=1):
                title = ""
                body_lines: list[str] = []
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text or "" for run in para.runs).strip()
                        if not text:
                            continue
                        if not title and shape == slide.shapes.title:
                            title = text
                        else:
                            body_lines.append(text)
                blocks.append(f"## Slide {i}" + (f": {title}" if title else ""))
                for line in body_lines:
                    blocks.append(f"- {line}")
                blocks.append("")
            out = "\n".join(blocks).strip()
            return out or None
    except Exception as exc:  # noqa: BLE001 — extraction is best-effort
        logger.warning("override text extraction failed path=%s mime=%s: %s", path, mime, exc)
        return None
    return None


def _override_content(row) -> str | None:
    """Return the wire `content` for an override upload.

    Priority:
      1. Textual mime (yaml/html/xml/…) → UTF-8 decode the file bytes.
      2. Binary mime we can extract text from (docx/xlsx/pptx) → the
         extracted markdown, so the partner preview can render inline
         alongside the original binary in the download slot.
      3. Anything else (pdf/mp4/zip/…) → the placeholder sentinel; the
         partner UI renders the "Binary attachment — Download to view"
         card and the Download button serves the pristine binary.

    Returns None when no override is set on this row."""
    path = getattr(row, "override_path", None)
    if not path or not os.path.exists(path):
        return None
    mime = (getattr(row, "override_mime_type", None) or "").lower()
    if mime in _TEXTUAL_OVERRIDE_MIMES:
        try:
            with open(path, "rb") as f:
                return f.read().decode("utf-8", errors="replace")
        except OSError:
            return _OVERRIDE_PLACEHOLDER
    extracted = _extract_text_from_upload(path, mime)
    if extracted:
        return extracted
    return _OVERRIDE_PLACEHOLDER


def _override_attach(out: dict, row, primary_kind: str, doc_type: str = "") -> bool:
    """Overwrite the row's primary attachment in the payload dict with the
    override upload. `primary_kind` names the attachment slot for this doc_type
    ("docx" / "pptx" / "video" / "xlsx" / "xsd_zip"). Returns True when the
    override was applied so the caller can skip its own generated-file read.

    Content-primary short-circuit: for doc_types the partner UI downloads via
    `/download/native` (manifest, prototype_screens, single-file xsd), when the
    override is textual (YAML / HTML / XML / JSON / plain / markdown), the
    upload's bytes are already written into the wire `content` field by the
    caller's `_override_content` — populating a binary slot on top would ship
    a mismatched `docx_b64` payload with the ACTUAL text sitting unused on
    the partner side. Skip the binary attach in that case; content wins.
    """
    path = getattr(row, "override_path", None)
    if not path or not os.path.exists(path):
        return False
    mime = (getattr(row, "override_mime_type", None) or "").lower()
    if _uses_content_primary(doc_type) and mime in _TEXTUAL_OVERRIDE_MIMES:
        # Textual override for a content-primary doc_type: content path only.
        # Return True so the caller skips its generated-file binary read; the
        # content field already carries the override bytes.
        return True
    try:
        with open(path, "rb") as f:
            raw = f.read()
    except OSError as exc:
        logger.warning("Skipping override attach for doc=%s: %s", getattr(row, "id", "?"), exc)
        return False
    # Drop any stale attachment fields for other kinds so a single unified
    # override on the wire replaces the whole generated-attachment set.
    for kind in ("docx", "pptx", "video", "xlsx", "xsd_zip"):
        for suffix in ("_b64", "_filename", "_sha256", "_size_bytes", "_mime_type", "_file_count"):
            out.pop(f"{kind}{suffix}", None)
    key = primary_kind
    out[f"{key}_b64"]        = base64.b64encode(raw).decode()
    out[f"{key}_filename"]   = getattr(row, "override_filename", None) or os.path.basename(path)
    out[f"{key}_sha256"]     = getattr(row, "override_sha256", None) \
        or hashlib.sha256(raw).hexdigest()
    out[f"{key}_size_bytes"] = getattr(row, "override_size_bytes", None) or len(raw)
    out[f"{key}_mime_type"]  = getattr(row, "override_mime_type", None) \
        or "application/octet-stream"
    return True


def _primary_kind_for_doc_type(dt: str) -> str:
    """Map a doc_type to the attachment slot the override replaces. When the
    generated artifact is a Word doc (BRD/product_note/etc.), the override lands
    in `docx_*`; product_deck → `pptx_*`; promo/explainer video → `video_*`;
    cert_test_cases → `xlsx_*`; xsd → `xsd_zip_*`."""
    if dt in ("promo_video", "explainer_video"):
        return "video"
    if dt == "product_deck":
        return "pptx"
    if dt == "cert_test_cases":
        return "xlsx"
    if dt == "xsd":
        return "xsd_zip"
    return "docx"


def build_kit_envelope(
    cr: ChangeRequest,
    db: Session,
    *,
    change_summary: str = "",
    include_doc_types: set[str] | None = None,
) -> dict:
    """Assemble the `product_kit` A2A payload from the latest kit docs + XSD.

    Pure packaging — no DB writes, no awaits. Latest version per doc_type only
    (ProductKitDocument keeps history), so a revision never ships duplicate
    doc_types.

    `change_summary` (Slice 4) is the partner-facing "what changed since the
    previous version" note attached on a revision ship; empty on the initial
    ship (v1).

    `include_doc_types` (migration 0115) — when supplied, doc types not in the
    set are skipped entirely from `documents[]`. When None, ship every eligible
    doc (legacy behaviour — keeps programmatic callers unchanged).
    """
    change_id = cr.id
    # Partner-facing kit version — all kit artifacts (incl. the XSD zip filename)
    # are labelled with this so they read consistently as "kit v{nv}".
    nv = getattr(cr, "negotiation_version", 1) or 1
    kit_docs = latest_kit_docs(db, change_id)

    # Resolve the cert_test_cases xlsx path once. The xlsx lives only on the
    # agent job's result_payload, not on the ProductKitDocument row.
    cert_xlsx_path: str | None = None
    try:
        from app.models.agent_job import AgentJob, AgentJobStatus
        cert_job = (
            db.query(AgentJob)
            .filter(
                AgentJob.change_request_id == change_id,
                AgentJob.module == "product_kit",
                AgentJob.subtype == "cert_test_cases",
                AgentJob.status == AgentJobStatus.SUCCEEDED,
            )
            .order_by(AgentJob.completed_at.desc().nullslast(), AgentJob.updated_at.desc())
            .first()
        )
        if cert_job is not None:
            rp = cert_job.result_payload or {}
            files = rp.get("files") or {}
            cert_xlsx_path = files.get("xlsx") or rp.get("xlsx_path")
    except Exception as exc:  # noqa: BLE001
        logger.warning("cert_test_cases xlsx lookup failed for change=%s: %s", change_id, exc)

    def _doc_payload(doc: ProductKitDocument) -> dict:
        dt_key = doc.doc_type.value if hasattr(doc.doc_type, 'value') else doc.doc_type
        override_text = _override_content(doc)
        effective_content = override_text if override_text is not None else (doc.content or "")
        out = {
            "doc_type": dt_key,
            "content":  effective_content,
            "version":  doc.version,
            "content_sha256": hashlib.sha256((effective_content or "").encode("utf-8")).hexdigest(),
        }
        if _override_attach(out, doc, _primary_kind_for_doc_type(dt_key), doc_type=dt_key):
            return out
        if doc.docx_path and os.path.exists(doc.docx_path):
            try:
                with open(doc.docx_path, "rb") as f:
                    raw = f.read()
                out["docx_b64"]       = base64.b64encode(raw).decode()
                out["docx_filename"]  = os.path.basename(doc.docx_path)
                out["docx_sha256"]    = hashlib.sha256(raw).hexdigest()
                out["docx_size_bytes"] = len(raw)
                out["docx_mime_type"] = (
                    "application/vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                )
            except OSError as exc:
                logger.warning("Skipping docx attach for doc=%s change=%s: %s", doc.id, change_id, exc)
        pptx_path = getattr(doc, "pptx_path", None)
        if pptx_path and os.path.exists(pptx_path):
            try:
                with open(pptx_path, "rb") as f:
                    out["pptx_b64"]      = base64.b64encode(f.read()).decode()
                    out["pptx_filename"] = os.path.basename(pptx_path)
            except OSError as exc:
                logger.warning("Skipping pptx attach for doc=%s change=%s: %s", doc.id, change_id, exc)
        # Promo/explainer video — the PM-uploaded MP4 lives at file_path; ship it
        # base64-inline like the other attachments (25 MB cap enforced at upload).
        dt0 = doc.doc_type.value if hasattr(doc.doc_type, 'value') else doc.doc_type
        video_path = getattr(doc, "file_path", None)
        if dt0 in ("promo_video", "explainer_video") and video_path and os.path.exists(video_path):
            try:
                with open(video_path, "rb") as f:
                    raw = f.read()
                out["video_b64"]        = base64.b64encode(raw).decode()
                out["video_filename"]   = os.path.basename(video_path)
                out["video_sha256"]     = hashlib.sha256(raw).hexdigest()
                out["video_size_bytes"] = len(raw)
                out["video_mime_type"]  = "video/mp4"
            except OSError as exc:
                logger.warning("Skipping video attach for doc=%s change=%s: %s", doc.id, change_id, exc)
        dt = doc.doc_type.value if hasattr(doc.doc_type, 'value') else doc.doc_type
        if dt == "cert_test_cases" and cert_xlsx_path and os.path.exists(cert_xlsx_path):
            try:
                with open(cert_xlsx_path, "rb") as f:
                    raw = f.read()
                out["xlsx_b64"]        = base64.b64encode(raw).decode()
                out["xlsx_filename"]   = os.path.basename(cert_xlsx_path)
                out["xlsx_sha256"]     = hashlib.sha256(raw).hexdigest()
                out["xlsx_size_bytes"] = len(raw)
                out["xlsx_mime_type"]  = (
                    "application/vnd.openxmlformats-officedocument."
                    "spreadsheetml.sheet"
                )
            except OSError as exc:
                logger.warning("Skipping xlsx attach for doc=%s change=%s: %s", doc.id, change_id, exc)
        return out

    # XSD ships alongside the kit whenever the Authority has produced one — same
    # documents[] envelope (doc_type="xsd"). Latest version, any non-empty
    # content. We ship REQUIRED schemas AND NOT-REQUIRED assessments so the
    # partner always sees the XSD agent's verdict for the change.
    #
    # NOTE: XSD/TSD are not negotiation-versioned (no negotiation_version column
    # like ProductKitDocument), so they ship as the latest working version, not
    # pinned to this kit's nv. The version that actually shipped IS recorded in
    # KitPublication.source_doc_versions (xsd/tsd are in documents[] below), so
    # the publication remains an exact, auditable snapshot. A true per-nv pin
    # would require adding negotiation_version to TechSpec/XSD — tracked, not done
    # here. In practice a revision regenerates all three together, so the window
    # for drift is a re-send of the same nv after an out-of-band XSD/TSD regen.
    xsd_row = (
        db.query(XSD)
        .filter(XSD.change_request_id == change_id)
        .order_by(XSD.version.desc())
        .first()
    )

    def _xsd_schema_zip(content: str) -> tuple[bytes, int] | None:
        """Bundle the change's `.xsd` schema files into one zip — the SAME set
        the Authority operator download serves, extracted from `XSD.content`'s fenced
        blocks via the shared helper (dedup logic mirrors that download).

        Returns None unless there are ≥2 schemas: a single-schema change keeps
        the partner's clean native `.xsd` download. The native download only
        returns the FIRST fenced block, so a multi-schema change needs this zip
        to deliver every file.
        """
        from app.services.xsd_bundle import extract_xsd_blocks
        blocks = extract_xsd_blocks(content or "")
        if len(blocks) < 2:
            return None
        buf = io.BytesIO()
        seen: dict[str, int] = {}
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for fname, body in blocks:
                # Disambiguate if the agent re-used a filename across blocks.
                if fname in seen:
                    seen[fname] += 1
                    stem, _, ext = fname.rpartition(".")
                    fname = f"{stem}_{seen[fname]}.{ext}" if ext else f"{fname}_{seen[fname]}"
                else:
                    seen[fname] = 1
                zf.writestr(fname, body)
        return buf.getvalue(), len(blocks)

    def _xsd_payload(row: XSD) -> dict:
        override_text = _override_content(row)
        effective_content = override_text if override_text is not None else (row.content or "")
        out = {
            "doc_type":       "xsd",
            "content":        effective_content,
            "version":        row.version,
            "content_sha256": hashlib.sha256((effective_content or "").encode("utf-8")).hexdigest(),
            "is_required":    row.is_required,
        }
        if _override_attach(out, row, "xsd_zip", doc_type="xsd"):
            return out
        zipped = _xsd_schema_zip(row.content)
        if zipped is not None:
            raw, count = zipped
            out["xsd_zip_b64"]        = base64.b64encode(raw).decode()
            out["xsd_zip_filename"]   = f"xsd_schemas_v{nv}.zip"
            out["xsd_zip_sha256"]     = hashlib.sha256(raw).hexdigest()
            out["xsd_zip_size_bytes"] = len(raw)
            out["xsd_zip_mime_type"]  = "application/zip"
            out["xsd_file_count"]     = count
        if row.docx_path and os.path.exists(row.docx_path):
            try:
                with open(row.docx_path, "rb") as f:
                    raw = f.read()
                out["docx_b64"]        = base64.b64encode(raw).decode()
                out["docx_filename"]   = os.path.basename(row.docx_path)
                out["docx_sha256"]     = hashlib.sha256(raw).hexdigest()
                out["docx_size_bytes"] = len(raw)
                out["docx_mime_type"]  = (
                    "application/vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                )
            except OSError as exc:
                logger.warning("Skipping xsd docx attach for change=%s: %s", change_id, exc)
        return out

    # TSD (Technical Spec Document) ships alongside the kit the same way XSD
    # does — same documents[] envelope (doc_type="tsd"), latest version with
    # non-empty content. Partners need the tech spec to plan their integration.
    tsd_row = (
        db.query(TechSpec)
        .filter(TechSpec.change_request_id == change_id)
        .order_by(TechSpec.version.desc())
        .first()
    )

    def _tsd_payload(row: TechSpec) -> dict:
        override_text = _override_content(row)
        effective_content = override_text if override_text is not None else (row.content or "")
        out = {
            "doc_type":       "tsd",
            "content":        effective_content,
            "version":        row.version,
            "content_sha256": hashlib.sha256((effective_content or "").encode("utf-8")).hexdigest(),
        }
        if _override_attach(out, row, "docx", doc_type="tsd"):
            return out
        if row.docx_path and os.path.exists(row.docx_path):
            try:
                with open(row.docx_path, "rb") as f:
                    raw = f.read()
                out["docx_b64"]        = base64.b64encode(raw).decode()
                out["docx_filename"]   = os.path.basename(row.docx_path)
                out["docx_sha256"]     = hashlib.sha256(raw).hexdigest()
                out["docx_size_bytes"] = len(raw)
                out["docx_mime_type"]  = (
                    "application/vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                )
            except OSError as exc:
                logger.warning("Skipping tsd docx attach for change=%s: %s", change_id, exc)
        return out

    # Promo/explainer videos: the SCRIPT is an internal the Authority artifact (the basis
    # for producing the video) and is NOT shipped to the partner. Ship the doc
    # only when an actual MP4 was uploaded (generated OR override), and only the
    # video (no script body).
    _VIDEO_TYPES = ("promo_video", "explainer_video")

    def _keep(dt: str) -> bool:
        return include_doc_types is None or dt in include_doc_types

    documents = []
    for doc in kit_docs:
        dt = doc.doc_type.value if hasattr(doc.doc_type, 'value') else doc.doc_type
        if not _keep(dt):
            continue
        has_override = bool(getattr(doc, "override_path", None) and os.path.exists(doc.override_path))
        if dt in _VIDEO_TYPES:
            if (doc.file_path and os.path.exists(doc.file_path)) or has_override:
                p = _doc_payload(doc)
                if not has_override:
                    p["content"] = ""
                    p["content_sha256"] = hashlib.sha256(b"").hexdigest()
                documents.append(p)
        elif doc.content or has_override:
            documents.append(_doc_payload(doc))
    if xsd_row and _keep("xsd") and (
        (xsd_row.content or "").strip()
        or (getattr(xsd_row, "override_path", None) and os.path.exists(xsd_row.override_path))
    ):
        documents.append(_xsd_payload(xsd_row))
    if tsd_row and _keep("tsd") and (
        (tsd_row.content or "").strip()
        or (getattr(tsd_row, "override_path", None) and os.path.exists(tsd_row.override_path))
    ):
        documents.append(_tsd_payload(tsd_row))

    return {
        "schema_version": "1.0",
        "message_kind": "CHANGE_COMMUNICATION",
        "kit_id": f"CHG_{change_id}",
        "kit_name": cr.title,
        "rollout_type": "STANDARD",
        "valid_until": (utcnow() + timedelta(days=7)).isoformat(),
        "change_id": change_id,
        # Partner-facing published version. The partner stores docs tagged with
        # this and bumps its mirror when it grows.
        "negotiation_version": nv,
        "title": cr.title,
        "initial_prompt": cr.initial_prompt,
        "enhanced_prompt": cr.enhanced_prompt,
        # Partner-facing summary of what changed from the previous version
        # (Slice 4). Empty on the initial ship.
        "change_summary": change_summary or "",
        "documents": documents,
    }


def snapshot_publication(
    cr: ChangeRequest,
    envelope: dict,
    db: Session,
    *,
    revision_reason: str | None = None,
    resolver_action: str | None = None,
    published_by: str | None = None,
) -> KitPublication:
    """Upsert an immutable KitPublication for the change's current
    negotiation_version. Re-send of the same version updates the row
    (respects the (change, version) unique constraint)."""
    nv = getattr(cr, "negotiation_version", 1) or 1
    sha = hashlib.sha256(
        json.dumps(envelope, sort_keys=True, default=str).encode("utf-8")
    ).hexdigest()
    source_doc_versions = {d["doc_type"]: d.get("version", 1) for d in envelope.get("documents", [])}

    row = db.scalars(
        select(KitPublication).where(
            KitPublication.change_request_id == cr.id,
            KitPublication.negotiation_version == nv,
        )
    ).first()
    if row is None:
        row = KitPublication(
            id=generate_uuid(),
            change_request_id=cr.id,
            negotiation_version=nv,
        )
        db.add(row)
    row.envelope = envelope
    row.envelope_sha256 = sha
    row.source_doc_versions = source_doc_versions
    row.revision_reason = revision_reason
    row.resolver_action = resolver_action
    row.published_by = published_by
    row.published_at = utcnow()
    db.commit()
    db.refresh(row)
    logger.info(
        "KitPublication: change=%s v%d sha=%s docs=%d",
        cr.id, nv, sha[:12], len(source_doc_versions),
    )
    return row


def _partner_inline_limit(partner: PartnerAgent) -> int:
    """Per-partner cap (bytes) on a single inline attachment's base64/wire size.
    Falls back to the global default; 0 = no limit (inline everything)."""
    v = getattr(partner, "max_inline_attachment_bytes", None)
    if v is None:
        v = settings.partner_max_inline_attachment_bytes
    try:
        return max(0, int(v or 0))
    except (TypeError, ValueError):
        return 0


def _gate_envelope_for_partner(envelope: dict, limit_bytes: int) -> tuple[dict, list[str]]:
    """Return an envelope copy with inline attachments whose base64 (wire) size
    exceeds ``limit_bytes`` OMITTED — the `<kind>_b64` blob is dropped and a
    `<kind>_omitted`/`<kind>_omitted_reason` marker added; all other metadata
    (filename, sha256, size, mime) is preserved so the partner still sees what
    exists and can request it out of band.

    ``limit_bytes <= 0`` means no limit → the envelope is returned unchanged.
    The copy is shallow-per-document (attachment values are strings), so the
    shared/snapshotted envelope passed in is never mutated. Returns (envelope,
    list-of-omission-notes) for logging.
    """
    if limit_bytes <= 0:
        return envelope, []
    notes: list[str] = []
    gated = dict(envelope)
    new_docs: list[dict] = []
    for doc in envelope.get("documents", []):
        d = dict(doc)  # shallow copy — attachment values are strings/scalars
        for key in [k for k in d if k.endswith("_b64")]:
            b64 = d[key]
            if isinstance(b64, str) and len(b64) > limit_bytes:
                kind = key[:-4]  # strip "_b64"
                del d[key]
                d[f"{kind}_omitted"] = True
                d[f"{kind}_omitted_reason"] = (
                    f"attachment exceeds partner inline limit "
                    f"({len(b64)} > {limit_bytes} wire bytes)"
                )
                notes.append(f"{d.get('doc_type', '?')}/{kind} ({len(b64)}B)")
        new_docs.append(d)
    gated["documents"] = new_docs
    return gated, notes


async def dispatch_kit_to_partners(
    cr: ChangeRequest,
    envelope: dict,
    assignments: list[ChangePartnerAssignment],
    db: Session,
    user_id: str | None,
    *,
    mode: str = "initial",
) -> dict:
    """Send the kit envelope to partners. See module docstring for `mode`.

    Attachments are size-gated per partner (`_gate_envelope_for_partner`) so a
    large inline video/docx can't exceed a partner's ingress body-size limit."""
    change_id = cr.id
    # Resolved once per dispatch, not per partner. A pack with no channel means
    # the domain has no way to reach implementers at all — the OCPP shape, where
    # a spec is published and read. Failing loudly here is right: the caller
    # asked to dispatch, and silently delivering nothing would leave every
    # assignment stuck with no signal.
    from app.core.domain.contract import (
        OutboundMessage as _OutboundMessage,
        Partner as _ChannelPartner,
        channel_of,
    )
    from app.core.domain.registry import get_active_pack

    _channel = channel_of(get_active_pack())
    # The channel manages its own persistence. Reaching in to hand it this
    # function's session would couple generic dispatch to one adapter's
    # internals to save a connection — the wrong trade. `send_task_to_partner`
    # already commits the audit row itself, so the transaction boundary is
    # unchanged either way.
    if _channel is None:
        raise RuntimeError(
            "kit dispatch requested but the active domain pack declares no "
            "partner channel; this domain publishes rather than dispatches"
        )

    results: list[dict] = []
    skipped = 0
    for assignment in assignments:
        if mode == "initial":
            # First send: only partners awaiting delivery; advance to RECEIVED.
            if assignment.status != AssignmentStatus.ASSIGNED:
                skipped += 1
                continue
        else:
            # Revision: re-send to every still-active partner, leaving the
            # assignment lifecycle untouched (gate is negotiation_version_accepted).
            if assignment.status == AssignmentStatus.WITHDRAWN:
                skipped += 1
                continue

        partner = db.get(PartnerAgent, assignment.partner_id)
        if not partner:
            continue

        # Size-gate attachments for THIS partner before shipping (per-partner
        # limit → global default → 0/no-limit).
        send_envelope, omitted = _gate_envelope_for_partner(
            envelope, _partner_inline_limit(partner)
        )
        if omitted:
            logger.info(
                "Kit dispatch: omitted %d oversize attachment(s) for partner=%s "
                "change=%s (limit=%dB): %s",
                len(omitted), partner.name, change_id,
                _partner_inline_limit(partner), ", ".join(omitted),
            )

        # Delivered through the active domain pack's channel rather than
        # calling A2A directly. Shipping a change to the parties who must
        # implement it is the platform's central act and the least
        # transport-specific thing it does — an ecosystem without an agent
        # protocol publishes instead, and that is a channel swap, not a
        # different feature.
        message = await _channel.deliver(
            _ChannelPartner(key=partner.id, label=partner.name or partner.id),
            _OutboundMessage(
                kind=A2ATaskType.CHANGE_COMMUNICATION.value,
                change_id=change_id,
                # `send_envelope` is already the wire payload, including any
                # inline base64 attachments left after per-partner size gating
                # (see _gate_envelope_for_partner). It is passed as payload
                # rather than OutboundMessage.attachments because the gating has
                # already made the transport's representation decision.
                payload=send_envelope,
            ),
        )

        # A failed send must NOT advance the partner's state. Previously RECEIVED was set
        # and the 24h round clock started even when `message.status == 'delivery_failed'`
        # — so a bank that never got the kit was recorded as having received it, and its
        # negotiation window silently expired. Only progress on an actual delivery.
        delivered = message.delivered
        if not delivered:
            logger.error(
                "kit dispatch NOT delivered — leaving assignment at %s and NOT starting the "
                "round clock: change=%s partner=%s status=%s error=%s task_id=%s",
                assignment.status, change_id, partner.name, message.status,
                message.error_code, message.reference,
            )
            # Delivery carries `status` and `error_code`, which is all this
            # reads off the message besides a partner-name fallback that the
            # PartnerAgent row already satisfies.
            notify_delivery_failure(db, message, partner, context="Product Kit dispatch")

        if mode == "initial" and delivered:
            set_status(
                assignment, AssignmentStatus.RECEIVED, db,
                actor_user_id=user_id,
                reason="Product Kit dispatched",
            )
            db.commit()

        # Start the 24h negotiation-round clock on delivery (Slice 2). The
        # round number tracks the published version: v1 → round 1, v2 →
        # round 2. v3 is the final kit, so no further round opens — the
        # negotiation freezes after round 2 (Slice 5). create_round_state is
        # idempotent, so a later PROPOSAL_ACKNOWLEDGED won't double-create.
        # Gated on delivery: no delivery, no clock.
        if delivered:
            from app.services.negotiation_extended import (
                MAX_ROUNDS,
                create_round_state,
                send_round_opened,
            )
            round_number = cr.negotiation_version or 1
            if round_number <= MAX_ROUNDS:
                _state, _was_created = create_round_state(
                    change_id, partner.id, round_number, db,
                )
                db.commit()
                # Only notify on a fresh round — a redispatch of the same
                # version hits the idempotent branch and mustn't re-signal.
                # Reason picks initial_ack for v1, version_ship for v2+; the
                # ack path in authority_handlers still emits initial_ack once the
                # partner PROPOSAL_ACKNOWLEDGED lands (idempotent — same
                # round_number, won't re-fire).
                if _was_created:
                    _reason = "initial_ack" if round_number == 1 else "version_ship"
                    try:
                        await send_round_opened(
                            change_request_id=change_id,
                            partner_id=partner.id,
                            round_number=round_number,
                            opened_reason=_reason,
                            db=db,
                        )
                    except Exception:
                        logger.exception(
                            "send_round_opened(%s) failed for change=%s partner=%s round=%d",
                            _reason, change_id, partner.id, round_number,
                        )

        results.append({
            "partner_id": partner.id,
            "partner_name": partner.name,
            "task_id": message.reference,
            "delivery_status": message.status,
        })

    # Freeze once a version beyond the round cap has shipped (the real-revision
    # path: MAX_ROUNDS rounds then the final kit). The no-change PM advance
    # freezes separately via advance_round_no_change; this is the ship backstop.
    # After freeze, the executor rejects inbound queries / counter-proposals;
    # the partner's only channel is an EmergencyIssue.
    from app.services.negotiation_extended import MAX_ROUNDS, notify_partners_frozen
    if (cr.negotiation_version or 1) > MAX_ROUNDS and cr.negotiation_frozen_at is None:
        cr.negotiation_frozen_at = utcnow()
        db.commit()
        logger.info("Negotiation frozen on v%d ship: change=%s", cr.negotiation_version, change_id)
        await notify_partners_frozen(change_id, db)

    logger.info(
        "Kit dispatch (%s): change=%s delivered=%d skipped=%d",
        mode, change_id, len(results), skipped,
    )
    return {"partners_notified": len(results), "partners_skipped": skipped, "results": results}
