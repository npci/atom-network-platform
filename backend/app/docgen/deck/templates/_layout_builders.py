# Copyright 2026 National Payments Corporation of India
# SPDX-License-Identifier: MIT

"""Pure shape-drawing helpers that emit one slide per call.

These are the bland-v1 visual treatments for each `SlideLayout` enum
value. Each builder accepts a `pptx.slide.Slide` (already inserted
into a presentation) plus the semantic content it needs, and adds
shapes/text frames in place. No I/O, no LLM, no schema dependency
— pure positioning + styling so the renderer (D4) and the master
generator can call the same code.

Visual contract (16:9, 13.33 × 7.5 inches):

    ┌──────────────────────────────────────────────┐
    │ TITLE BAR (0.9" tall, dark grey, white text) │  ← every layout except TITLE/SECTION
    ├──────────────────────────────────────────────┤
    │                                              │
    │  body region                                 │
    │  (varies per layout)                         │
    │                                              │
    ├──────────────────────────────────────────────┤
    │  footer: slide N / total                     │  ← every layout
    └──────────────────────────────────────────────┘

All shapes are real `Shape` / `TextBox` instances with editable text
frames — not flattened text — so the user can click into PowerPoint
and edit any element.
"""
from __future__ import annotations

from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationT
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt


# ── 16:9 canvas ──────────────────────────────────────────────────────────────

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN       = Inches(0.5)
TITLE_BAR_H  = Inches(0.9)
FOOTER_H     = Inches(0.35)
BODY_TOP     = TITLE_BAR_H + Inches(0.2)
BODY_BOTTOM  = SLIDE_HEIGHT - FOOTER_H - Inches(0.1)


# ── colour palette (bland v1 — designer replaces) ────────────────────────────

DARK    = RGBColor(0x40, 0x40, 0x40)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT  = RGBColor(0x1F, 0x4E, 0x79)   # neutral steel blue
LIGHT   = RGBColor(0xF2, 0xF2, 0xF2)
MID     = RGBColor(0xBF, 0xBF, 0xBF)
TEXT_PRIMARY   = RGBColor(0x20, 0x20, 0x20)
TEXT_SECONDARY = RGBColor(0x55, 0x55, 0x55)


# ── presentation factory ─────────────────────────────────────────────────────

def new_presentation() -> PresentationT:
    """Fresh blank 16:9 presentation. All builders assume this canvas."""
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ── shared building blocks ───────────────────────────────────────────────────

_BLANK_LAYOUT_INDEX = 6   # python-pptx default master: layout 6 is "Blank"


def _add_blank_slide(prs: PresentationT) -> Slide:
    return prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT_INDEX])


def _solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    # Table cells don't expose a line attribute; only suppress borders
    # for shapes that have one.
    line = getattr(shape, "line", None)
    if line is not None:
        line.fill.background()


def _set_text(
    text_frame,
    text: str,
    *,
    size_pt: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_PRIMARY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.vertical_anchor = anchor
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_bar(slide: Slide, title: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, TITLE_BAR_H,
    )
    _solid_fill(bar, DARK)
    bar.text_frame.margin_left = Inches(0.5)
    bar.text_frame.margin_right = Inches(0.5)
    _set_text(
        bar.text_frame, title,
        size_pt=28, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_footer(slide: Slide, slide_no: int, total: int, feature_name: str = "") -> None:
    # Left side: feature name (small, muted)
    if feature_name:
        left = slide.shapes.add_textbox(
            MARGIN, SLIDE_HEIGHT - FOOTER_H, Inches(8.0), FOOTER_H,
        )
        _set_text(
            left.text_frame, feature_name,
            size_pt=9, color=TEXT_SECONDARY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    # Right side: slide number
    right = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(2.0) - MARGIN, SLIDE_HEIGHT - FOOTER_H,
        Inches(2.0), FOOTER_H,
    )
    _set_text(
        right.text_frame, f"{slide_no} / {total}",
        size_pt=9, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _set_speaker_notes(slide: Slide, notes: str) -> None:
    if not notes:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes


# ── per-layout builders ──────────────────────────────────────────────────────

def build_title_slide(
    prs: PresentationT,
    *,
    title: str,
    subtitle: Optional[str],
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)

    # Centred big title at vertical 1/3
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(2.4), SLIDE_WIDTH - 2 * MARGIN, Inches(1.6),
    )
    _set_text(
        title_box.text_frame, title,
        size_pt=48, bold=True, color=DARK, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Subtitle below
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            MARGIN, Inches(4.2), SLIDE_WIDTH - 2 * MARGIN, Inches(0.9),
        )
        _set_text(
            sub_box.text_frame, subtitle,
            size_pt=20, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Decorative accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SLIDE_WIDTH / 2 - Inches(1.0), Inches(5.4),
        Inches(2.0), Emu(38100),  # 0.04"
    )
    _solid_fill(bar, ACCENT)

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_section_slide(
    prs: PresentationT,
    *,
    title: str,
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)

    # Full-bleed dark band fills middle 60% of the slide
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(2.0), SLIDE_WIDTH, Inches(3.5),
    )
    _solid_fill(band, ACCENT)

    headline = slide.shapes.add_textbox(
        MARGIN, Inches(2.0), SLIDE_WIDTH - 2 * MARGIN, Inches(3.5),
    )
    _set_text(
        headline.text_frame, title,
        size_pt=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_bullet_list_slide(
    prs: PresentationT,
    *,
    title: str,
    bullets: list[str],
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    body = slide.shapes.add_textbox(
        MARGIN + Inches(0.3), BODY_TOP + Inches(0.2),
        SLIDE_WIDTH - 2 * MARGIN - Inches(0.3),
        BODY_BOTTOM - BODY_TOP - Inches(0.4),
    )
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.size = Pt(20)
            p.runs[0].font.color.rgb = TEXT_PRIMARY

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def _add_column_block(
    slide: Slide,
    *,
    left: int, top: int, width: int, height: int,
    icon_hint: Optional[str], heading: str, body: str,
) -> None:
    # Outline rectangle for the column
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = MID
    box.line.width = Emu(6350)  # 0.5pt
    # Box has its own text frame; clear it so we can add stacked content
    box.text_frame.clear()
    box.text_frame.text = ""

    inner_pad = Inches(0.25)

    # Optional icon glyph (placeholder — bland v1 uses initial of icon_hint
    # in a circle; designer replaces with proper iconography later)
    icon_h = Inches(0.7)
    if icon_hint:
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + width // 2 - icon_h // 2, top + inner_pad,
            icon_h, icon_h,
        )
        _solid_fill(icon, ACCENT)
        _set_text(
            icon.text_frame, (icon_hint[:1] or "•").upper(),
            size_pt=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        heading_top = top + inner_pad + icon_h + Inches(0.15)
    else:
        heading_top = top + inner_pad

    # Heading
    head_box = slide.shapes.add_textbox(
        left + inner_pad, heading_top,
        width - 2 * inner_pad, Inches(0.5),
    )
    _set_text(
        head_box.text_frame, heading,
        size_pt=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )

    # Body
    body_top = heading_top + Inches(0.6)
    body_box = slide.shapes.add_textbox(
        left + inner_pad, body_top,
        width - 2 * inner_pad, top + height - body_top - inner_pad,
    )
    _set_text(
        body_box.text_frame, body,
        size_pt=12, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )


def build_two_column_slide(
    prs: PresentationT,
    *,
    title: str,
    columns,                  # list of 2 ColumnBlock (heading, body, icon_hint)
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    avail_w = SLIDE_WIDTH - 2 * MARGIN
    gap = Inches(0.3)
    col_w = (avail_w - gap) // 2
    body_h = BODY_BOTTOM - BODY_TOP - Inches(0.2)

    for i, c in enumerate(columns):
        left = MARGIN + i * (col_w + gap)
        _add_column_block(
            slide,
            left=left, top=BODY_TOP + Inches(0.1),
            width=col_w, height=body_h,
            icon_hint=getattr(c, "icon_hint", None),
            heading=c.heading,
            body=c.body,
        )

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_three_column_slide(
    prs: PresentationT,
    *,
    title: str,
    columns,                  # list of 3 ColumnBlock
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    avail_w = SLIDE_WIDTH - 2 * MARGIN
    gap = Inches(0.3)
    col_w = (avail_w - 2 * gap) // 3
    body_h = BODY_BOTTOM - BODY_TOP - Inches(0.2)

    for i, c in enumerate(columns):
        left = MARGIN + i * (col_w + gap)
        _add_column_block(
            slide,
            left=left, top=BODY_TOP + Inches(0.1),
            width=col_w, height=body_h,
            icon_hint=getattr(c, "icon_hint", None),
            heading=c.heading,
            body=c.body,
        )

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_numbered_flow_slide(
    prs: PresentationT,
    *,
    title: str,
    steps,                    # list of NumberedStep (n, label, body?)
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    n_steps = len(steps)
    avail_w = SLIDE_WIDTH - 2 * MARGIN
    gap = Inches(0.2)
    box_w = (avail_w - (n_steps - 1) * gap) // n_steps
    box_h = Inches(2.8)
    box_top = BODY_TOP + Inches(0.6)

    arrow_w = Inches(0.18)

    for i, step in enumerate(steps):
        left = MARGIN + i * (box_w + gap)
        # Step box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, box_w, box_h,
        )
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = ACCENT
        box.line.width = Emu(12700)  # 1pt
        box.text_frame.clear(); box.text_frame.text = ""

        # Number circle at top
        num_d = Inches(0.65)
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + box_w // 2 - num_d // 2, box_top + Inches(0.25),
            num_d, num_d,
        )
        _solid_fill(num, ACCENT)
        _set_text(
            num.text_frame, str(step.n),
            size_pt=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label
        lbl = slide.shapes.add_textbox(
            left + Inches(0.15), box_top + Inches(1.05),
            box_w - Inches(0.3), Inches(0.5),
        )
        _set_text(
            lbl.text_frame, step.label,
            size_pt=14, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )

        # Body (optional)
        if getattr(step, "body", None):
            bd = slide.shapes.add_textbox(
                left + Inches(0.15), box_top + Inches(1.6),
                box_w - Inches(0.3), Inches(1.1),
            )
            _set_text(
                bd.text_frame, step.body,
                size_pt=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.TOP,
            )

        # Right-arrow between this and next step
        if i < n_steps - 1:
            arrow_left = left + box_w + (gap - arrow_w) // 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left, box_top + box_h // 2 - Inches(0.18),
                arrow_w, Inches(0.36),
            )
            _solid_fill(arrow, ACCENT)

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_diagram_slide(
    prs: PresentationT,
    *,
    title: str,
    image_bytes: Optional[bytes],
    fallback_text: Optional[str],
    caption: Optional[str],
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    """Diagram slide. Either `image_bytes` (rendered PNG) OR
    `fallback_text` (the dot source as a code block, used when render
    fails). Caller decides which is available."""
    import io
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    region_top = BODY_TOP + Inches(0.2)
    caption_h = Inches(0.45) if caption else Emu(0)
    region_h = BODY_BOTTOM - region_top - caption_h - Inches(0.1)

    if image_bytes:
        from PIL import Image as PILImage
        # SCR #10 (Improper Resource Shutdown or Release): `Image.open` is lazy
        # and keeps a reference to the underlying buffer until closed. The
        # buffer here is in-memory, so this leaked no OS file descriptor, but
        # the object was never released deterministically and the scanner is
        # right that the acquire/release pair was unbalanced. Only `.size` is
        # needed, so scope it to a `with` and keep just the two integers.
        with PILImage.open(io.BytesIO(image_bytes)) as img:
            img_w_px, img_h_px = img.size
        img_aspect = img_w_px / max(img_h_px, 1)

        max_w = SLIDE_WIDTH - 2 * MARGIN - Inches(0.4)
        max_h = region_h
        max_aspect = max_w / max_h

        if img_aspect >= max_aspect:
            draw_w, draw_h = max_w, int(max_w / img_aspect)
        else:
            draw_h, draw_w = max_h, int(max_h * img_aspect)

        left = (SLIDE_WIDTH - draw_w) // 2
        top  = region_top + (region_h - draw_h) // 2
        slide.shapes.add_picture(
            io.BytesIO(image_bytes), left, top, width=draw_w, height=draw_h,
        )
    else:
        # Fallback: code block with the source
        fb = slide.shapes.add_textbox(
            MARGIN, region_top,
            SLIDE_WIDTH - 2 * MARGIN, region_h,
        )
        _set_text(
            fb.text_frame, fallback_text or "(diagram unavailable)",
            size_pt=10, color=TEXT_SECONDARY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
        # Switch font of all runs to monospace
        for p in fb.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Courier New"

    if caption:
        cap = slide.shapes.add_textbox(
            MARGIN, BODY_BOTTOM - caption_h - Inches(0.05),
            SLIDE_WIDTH - 2 * MARGIN, caption_h,
        )
        _set_text(
            cap.text_frame, caption,
            size_pt=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide


def build_table_slide(
    prs: PresentationT,
    *,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    slide_no: int,
    total: int,
    feature_name: str = "",
    speaker_notes: str = "",
) -> Slide:
    slide = _add_blank_slide(prs)
    _add_title_bar(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1   # +1 for header row

    table_left   = MARGIN
    table_top    = BODY_TOP + Inches(0.2)
    table_width  = SLIDE_WIDTH - 2 * MARGIN
    table_height = BODY_BOTTOM - table_top - Inches(0.2)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, table_left, table_top, table_width, table_height,
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        _solid_fill(cell, ACCENT)
        _set_text(
            cell.text_frame, h,
            size_pt=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        )
        cell.text_frame.margin_left = Inches(0.1)
        cell.text_frame.margin_top = Inches(0.05)

    # Body rows — alternate banding
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            _solid_fill(cell, LIGHT if i % 2 == 1 else WHITE)
            _set_text(
                cell.text_frame, v,
                size_pt=12, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT,
            )
            cell.text_frame.margin_left = Inches(0.1)
            cell.text_frame.margin_top = Inches(0.05)

    _add_footer(slide, slide_no, total, feature_name)
    _set_speaker_notes(slide, speaker_notes)
    return slide
